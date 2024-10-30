from click import prompt
import streamlit as st
from PyPDF2 import PdfReader
from docx import Document  # DOCX reader
from pptx import Presentation  # PPTX reader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

# Load API keys and environment variables
# Configuring the API keys 
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))



    
# Function to extract text from PDF
def extract_text_from_pdf(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

# Function to extract text from DOCX
def extract_text_from_docx(docx_docs):
    text = ""
    for docx in docx_docs:
        doc = Document(docx)
        for para in doc.paragraphs:
            text += para.text + "\n"
    return text

# Function to extract text from PPTX
def extract_text_from_pptx(pptx_docs):
    text = ""
    for pptx in pptx_docs:
        prs = Presentation(pptx)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

# Function to split text into chunks
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

# Function to get a vector store
def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

# Function to define the QA conversational chain
def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context. If the answer is not in the context, respond with "answer is not available in the context."\n\n
    Context:\n {context}?\n
    Question: \n{question}\n
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

# Function to handle user input and return a response
def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    st.write("Reply: ", response["output_text"])

def summarize_text(text):
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.5)
    prompt = f"Summarize the following text:\n\n{text[:1000]}" 
    response = model.predict(prompt)
    print(f"Model Response: {response}")  
    if isinstance(response, str):
        return response
    else:
        return response.get("text", "No summary available.")

# Main application
def main():
    print("Done")
    st.set_page_config("Chat with Documents")
    st.header("Chat with Your Documents (PDF, DOCX, PPTX)")

    # User input for questions
    user_question = st.text_input("Ask a Question from the Uploaded Files")

    if user_question:
        user_input(user_question)

    # Sidebar for file upload
    with st.sidebar:
        st.title("Menu:")
        pdf_docs = st.file_uploader("Upload PDF Files", accept_multiple_files=True, type=["pdf"])
        docx_docs = st.file_uploader("Upload DOCX Files", accept_multiple_files=True, type=["docx"])
        pptx_docs = st.file_uploader("Upload PPTX Files", accept_multiple_files=True, type=["pptx"])

        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                raw_text = ""

                # Process PDF Files
                if pdf_docs: 
                    raw_text += extract_text_from_pdf(pdf_docs)
                    summary = summarize_text(raw_text)
                    st.success("PDF Summary:")
                    st.write(summary)
                   

                # Process DOCX Files
                if docx_docs:
                    raw_text += extract_text_from_docx(docx_docs)

                # Process PPTX Files
                if pptx_docs:
                    raw_text += extract_text_from_pptx(pptx_docs)

                # Split text and build vector store
                if raw_text:
                    text_chunks = get_text_chunks(raw_text)
                    get_vector_store(text_chunks)
                    st.success("Files Processed Successfully")

if __name__ == "__main__":
    main()
